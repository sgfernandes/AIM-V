"""
AIM-V PowerPoint Slide Deck Generator
=======================================
Produces a 16:9 .pptx presentation with:
  - High-level AIM-V overview (architecture, agents, API)
  - Deep-dive on the Analytics Agent (OLS, QA/QC, cookbook compliance)

Usage:
    python generate_slides_pptx.py
    -> outputs: docs/AIM-V_Slides.pptx
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

os.makedirs("docs", exist_ok=True)

# ── Colour palette ───────────────────────────────────────────────────
BG       = RGBColor(0x0D, 0x1B, 0x2A)
BG_CARD  = RGBColor(0x1B, 0x28, 0x38)
ACCENT   = RGBColor(0x00, 0xB4, 0xD8)
ACCENT2  = RGBColor(0x00, 0x77, 0xB6)
WHITE    = RGBColor(0xE0, 0xE1, 0xDD)
GREEN    = RGBColor(0x06, 0xD6, 0xA0)
RED      = RGBColor(0xEF, 0x47, 0x6F)
ORANGE   = RGBColor(0xFF, 0xD1, 0x66)
MUTED    = RGBColor(0x77, 0x8D, 0xA9)

SW = Inches(13.333)  # 16:9 slide width
SH = Inches(7.5)

prs = Presentation()
prs.slide_width = SW
prs.slide_height = SH

BLANK_LAYOUT = prs.slide_layouts[6]  # blank


# ── Helpers ──────────────────────────────────────────────────────────
def set_slide_bg(slide, color=BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, *,
                font_size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT,
                font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_frame(slide, left, top, width, height, items, *,
                     font_size=14, color=WHITE, bullet_color=ACCENT,
                     spacing_pt=6):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = Pt(spacing_pt)
        p.level = 0
        # Bullet character
        pPr = p._pPr
        if pPr is None:
            from pptx.oxml.ns import qn
            pPr = p._p.get_or_add_pPr()
        from pptx.oxml.ns import qn
        from lxml import etree
        buNone = pPr.find(qn("a:buNone"))
        if buNone is not None:
            pPr.remove(buNone)
        buChar = etree.SubElement(pPr, qn("a:buChar"))
        buChar.set("char", "\u25b8")
        buClr = etree.SubElement(pPr, qn("a:buClr"))
        srgb = etree.SubElement(buClr, qn("a:srgbClr"))
        srgb.set("val", f"{bullet_color}")
        buSzPct = etree.SubElement(pPr, qn("a:buSzPct"))
        buSzPct.set("val", "100000")
    return txBox


def add_card(slide, left, top, width, height, title, body_lines,
             border_color=ACCENT):
    # Card background
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = BG_CARD
    shape.line.color.rgb = border_color
    shape.line.width = Pt(1.5)

    # Title
    add_textbox(slide, left + Inches(0.15), top + Inches(0.05),
                width - Inches(0.3), Inches(0.4), title,
                font_size=13, color=border_color, bold=True,
                alignment=PP_ALIGN.CENTER)

    # Separator line
    slide.shapes.add_connector(1,  # straight
        left + Inches(0.2), top + Inches(0.45),
        left + width - Inches(0.2), top + Inches(0.45)
    ).line.color.rgb = border_color

    # Body
    add_bullet_frame(slide, left + Inches(0.15), top + Inches(0.5),
                     width - Inches(0.3), height - Inches(0.6),
                     body_lines, font_size=10, spacing_pt=2)


def add_footer(slide, num, total):
    add_textbox(slide, Inches(0.4), SH - Inches(0.4), Inches(5), Inches(0.3),
                "AIM-V \u2014 Automation for Industrial M&V",
                font_size=8, color=MUTED)
    add_textbox(slide, SW - Inches(1.2), SH - Inches(0.4), Inches(0.8), Inches(0.3),
                f"{num}/{total}", font_size=8, color=MUTED,
                alignment=PP_ALIGN.RIGHT)


def add_divider(slide, left, top, width, color=ACCENT):
    line = slide.shapes.add_connector(1, left, top, left + width, top)
    line.line.color.rgb = color
    line.line.width = Pt(2)


TOTAL = 12

# ====================================================================
#  SLIDE 1 — Title
# ====================================================================
sl = prs.slides.add_slide(BLANK_LAYOUT)
set_slide_bg(sl)
add_textbox(sl, Inches(0), Inches(1.8), SW, Inches(1.5), "AIM-V",
            font_size=60, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
add_textbox(sl, Inches(0), Inches(3.2), SW, Inches(0.8),
            "Automation for Industrial M&V",
            font_size=28, color=ACCENT, alignment=PP_ALIGN.CENTER)
add_divider(sl, Inches(3.5), Inches(4.1), Inches(6.3))
add_textbox(sl, Inches(0), Inches(4.3), SW, Inches(0.6),
            "A Multi-Agent Platform for DOE ITV M&V Workflows",
            font_size=16, color=MUTED, alignment=PP_ALIGN.CENTER)
add_textbox(sl, Inches(0), Inches(5.2), SW, Inches(0.5),
            "Strategy  \u2022  Analytics  \u2022  Documentation",
            font_size=18, color=WHITE, alignment=PP_ALIGN.CENTER)
add_footer(sl, 1, TOTAL)

# ====================================================================
#  SLIDE 2 — The Challenge
# ====================================================================
sl = prs.slides.add_slide(BLANK_LAYOUT)
set_slide_bg(sl)
add_textbox(sl, Inches(0.5), Inches(0.3), Inches(6), Inches(0.7),
            "The Challenge", font_size=32, bold=True, color=WHITE)
add_divider(sl, Inches(0.5), Inches(1.0), Inches(4.5))

problems = [
    "Manual M&V processes are slow, error-prone, and expensive",
    "ASHRAE G14 / IPMVP compliance requires 15+ statistical checks",
    "Typical M&V plan takes weeks of analyst time per project",
    "No consistent QA/QC framework across an ESCO portfolio",
    "Reporting is fragmented \u2014 spreadsheets, PDFs, emails",
]
add_bullet_frame(sl, Inches(0.5), Inches(1.2), Inches(7.5), Inches(4),
                 problems, font_size=16, spacing_pt=10)

# Stat callout boxes
stats = [("15+", "QA/QC Checks"), ("3\u20135 wks", "Per Plan"), ("\u223c$50K+", "Analyst Cost / yr")]
for i, (val, label) in enumerate(stats):
    top = Inches(1.3 + i * 1.7)
    shape = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(9.2), top, Inches(3.2), Inches(1.3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = BG_CARD
    shape.line.color.rgb = ACCENT
    shape.line.width = Pt(1.5)
    add_textbox(sl, Inches(9.2), top + Inches(0.1), Inches(3.2), Inches(0.7),
                val, font_size=30, bold=True, color=ACCENT, alignment=PP_ALIGN.CENTER)
    add_textbox(sl, Inches(9.2), top + Inches(0.75), Inches(3.2), Inches(0.4),
                label, font_size=12, color=MUTED, alignment=PP_ALIGN.CENTER)

add_footer(sl, 2, TOTAL)

# ====================================================================
#  SLIDE 3 — AIM-V at a Glance
# ====================================================================
sl = prs.slides.add_slide(BLANK_LAYOUT)
set_slide_bg(sl)
add_textbox(sl, Inches(0.5), Inches(0.3), Inches(8), Inches(0.7),
            "AIM-V at a Glance", font_size=32, bold=True, color=WHITE)
add_divider(sl, Inches(0.5), Inches(1.0), Inches(5))

features = [
    "Multi-agent architecture \u2014 three specialized engines",
    "Conversational /chat API routes prompts to the right agent",
    "ITV-MV Cookbook & ASHRAE Guideline 14 compliant analytics",
    "Strategy engine infers IPMVP Option A/B/C/D automatically",
    "Documentation engine generates M&V plan markdown on demand",
    "Streamlit UI for interactive workflow exploration",
    "FastAPI backend with Swagger docs at /docs",
]
add_bullet_frame(sl, Inches(0.5), Inches(1.2), Inches(11), Inches(5),
                 features, font_size=16, spacing_pt=12)

add_footer(sl, 3, TOTAL)

# ====================================================================
#  SLIDE 4 — System Architecture
# ====================================================================
sl = prs.slides.add_slide(BLANK_LAYOUT)
set_slide_bg(sl)
add_textbox(sl, Inches(0), Inches(0.2), SW, Inches(0.7),
            "System Architecture", font_size=30, bold=True, color=WHITE,
            alignment=PP_ALIGN.CENTER)
add_divider(sl, Inches(3), Inches(0.9), Inches(7.3))

# Architecture boxes
arch_boxes = [
    (0.8,  3.0, 2.5, 1.0, "User\n(Chat / UI)",         MUTED),
    (4.2,  3.0, 3.0, 1.0, "Orchestrator\n(Intent Router)", ACCENT),
    (8.5,  1.2, 3.0, 0.9, "Strategy Engine",            ORANGE),
    (8.5,  3.0, 3.0, 0.9, "Analytics Engine",           GREEN),
    (8.5,  4.8, 3.0, 0.9, "Documentation Engine",       RED),
    (4.2,  5.5, 3.0, 0.8, "FastAPI Backend",            ACCENT2),
]
for bx, by, bw, bh, label, col in arch_boxes:
    shape = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(bx), Inches(by), Inches(bw), Inches(bh))
    shape.fill.solid()
    shape.fill.fore_color.rgb = BG_CARD
    shape.line.color.rgb = col
    shape.line.width = Pt(2)
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = col
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].space_before = Pt(4)

# Arrows (simple connectors)
arrows = [
    (3.3, 3.5, 4.2, 3.5),   # User -> Orch
    (7.2, 2.8, 8.5, 1.7),   # Orch -> Strategy
    (7.2, 3.5, 8.5, 3.5),   # Orch -> Analytics
    (7.2, 4.2, 8.5, 5.2),   # Orch -> Documentation
    (5.7, 4.0, 5.7, 5.5),   # Orch -> Backend
]
for x1, y1, x2, y2 in arrows:
    conn = sl.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    conn.line.color.rgb = MUTED
    conn.line.width = Pt(1.5)

# Phase labels
add_textbox(sl, Inches(11.6), Inches(1.2), Inches(1.5), Inches(0.3),
            "Phase 1", font_size=9, color=ORANGE)
add_textbox(sl, Inches(11.6), Inches(3.0), Inches(1.5), Inches(0.3),
            "Phase 2 & 4", font_size=9, color=GREEN)
add_textbox(sl, Inches(11.6), Inches(4.8), Inches(1.5), Inches(0.3),
            "Phase 3 & 5", font_size=9, color=RED)

add_footer(sl, 4, TOTAL)

# ====================================================================
#  SLIDE 5 — Three Agents
# ====================================================================
sl = prs.slides.add_slide(BLANK_LAYOUT)
set_slide_bg(sl)
add_textbox(sl, Inches(0), Inches(0.2), SW, Inches(0.7),
            "Three Specialized Agents", font_size=30, bold=True, color=WHITE,
            alignment=PP_ALIGN.CENTER)
add_divider(sl, Inches(3), Inches(0.9), Inches(7.3))

add_card(sl, Inches(0.3), Inches(1.3), Inches(3.9), Inches(5.5),
         "Strategy Engine", [
    "Infers IPMVP Option (A/B/C/D)",
    "Identifies independent variables",
    "Sets measurement boundary",
    "Defines assumptions",
    "Outputs next actions for plan setup",
], border_color=ORANGE)

add_card(sl, Inches(4.6), Inches(1.3), Inches(3.9), Inches(5.5),
         "Analytics Engine", [
    "OLS baseline regression",
    "15+ ITV-MV Cookbook QA/QC checks",
    "R\u00b2, CV(RMSE), NDBE, F-stat",
    "Autocorrelation, Shapiro-Wilk",
    "t-stat, p-values, outlier detection",
    "Post-period savings & FSU",
], border_color=GREEN)

add_card(sl, Inches(8.9), Inches(1.3), Inches(3.9), Inches(5.5),
         "Documentation Engine", [
    "Generates M&V plan markdown",
    "Integrates Strategy & Analytics",
    "Project / facility / boundary info",
    "Model stats & QA/QC summary",
    "Export-ready (PDF / Word / HTML)",
], border_color=RED)

add_footer(sl, 5, TOTAL)

# ====================================================================
#  SLIDE 6 — Analytics Deep Dive
# ====================================================================
sl = prs.slides.add_slide(BLANK_LAYOUT)
set_slide_bg(sl)
add_textbox(sl, Inches(0.5), Inches(0.3), Inches(10), Inches(0.7),
            "Analytics Agent \u2014 Deep Dive", font_size=32, bold=True, color=GREEN)
add_divider(sl, Inches(0.5), Inches(1.0), Inches(6.5), color=GREEN)

items = [
    "Core engine: statsmodels OLS regression on user-supplied baseline data",
    "Accepts any columns \u2014 auto-detects dependent variable & predictors",
    "Validates against ITV-MV Cookbook / ASHRAE Guideline 14 thresholds",
    "All thresholds are configurable via context overrides",
    "Returns structured JSON \u2014 ready for downstream agents or UI rendering",
    "Optional post-period savings & Fractional Savings Uncertainty (FSU)",
]
add_bullet_frame(sl, Inches(0.5), Inches(1.3), Inches(11), Inches(4.5),
                 items, font_size=16, spacing_pt=14)

# Tech badges
techs = [("statsmodels", GREEN), ("scipy", ACCENT), ("numpy", ORANGE),
         ("pandas", RED), ("FastAPI", ACCENT2)]
for i, (tech, col) in enumerate(techs):
    left = Inches(0.8 + i * 2.4)
    shape = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 left, Inches(6.0), Inches(1.8), Inches(0.45))
    shape.fill.solid()
    shape.fill.fore_color.rgb = BG_CARD
    shape.line.color.rgb = col
    shape.line.width = Pt(1.5)
    tf = shape.text_frame
    p = tf.paragraphs[0]
    p.text = tech
    p.font.size = Pt(11)
    p.font.color.rgb = col
    p.font.bold = True
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER

add_footer(sl, 6, TOTAL)

# ====================================================================
#  SLIDE 7 — Model-Level QA/QC
# ====================================================================
sl = prs.slides.add_slide(BLANK_LAYOUT)
set_slide_bg(sl)
add_textbox(sl, Inches(0.5), Inches(0.3), Inches(10), Inches(0.7),
            "Analytics \u2014 Model-Level QA/QC", font_size=30, bold=True, color=GREEN)
add_divider(sl, Inches(0.5), Inches(1.0), Inches(7), color=GREEN)

checks = [
    ("R-squared (R\u00b2)",          "\u2265 0.75",  "Goodness of fit"),
    ("CV(RMSE)",                    "\u2264 20%",   "Model prediction accuracy"),
    ("NDBE",                        "\u2264 0.5%",  "Net determination bias error"),
    ("F-statistic p-value",         "< 0.05",      "Overall model significance"),
    ("Autocorrelation |\u03c1|",    "< 0.70",      "Residual independence"),
]

# Table header background
hdr_shape = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE,
    Inches(0.5), Inches(1.3), Inches(12), Inches(0.5))
hdr_shape.fill.solid()
hdr_shape.fill.fore_color.rgb = ACCENT2
hdr_shape.line.fill.background()

add_textbox(sl, Inches(0.7), Inches(1.3), Inches(4), Inches(0.5),
            "Check", font_size=14, bold=True, color=WHITE)
add_textbox(sl, Inches(5.0), Inches(1.3), Inches(2), Inches(0.5),
            "Threshold", font_size=14, bold=True, color=WHITE)
add_textbox(sl, Inches(7.5), Inches(1.3), Inches(5), Inches(0.5),
            "Purpose", font_size=14, bold=True, color=WHITE)

for i, (check, thresh, purpose) in enumerate(checks):
    top = Inches(1.9 + i * 0.7)
    # Alternating row background
    if i % 2 == 0:
        row_bg = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE,
            Inches(0.5), top, Inches(12), Inches(0.6))
        row_bg.fill.solid()
        row_bg.fill.fore_color.rgb = BG_CARD
        row_bg.line.fill.background()

    add_textbox(sl, Inches(0.7), top + Inches(0.05), Inches(4), Inches(0.5),
                check, font_size=14, color=WHITE)
    add_textbox(sl, Inches(5.0), top + Inches(0.05), Inches(2), Inches(0.5),
                thresh, font_size=14, color=GREEN, bold=True)
    add_textbox(sl, Inches(7.5), top + Inches(0.05), Inches(5), Inches(0.5),
                purpose, font_size=13, color=MUTED)

add_textbox(sl, Inches(0.5), Inches(5.6), Inches(11), Inches(0.5),
            "All checks run automatically.  model_pass = True only when ALL pass.",
            font_size=14, color=ORANGE, bold=False)
add_textbox(sl, Inches(0.5), Inches(6.1), Inches(11), Inches(0.4),
            "ITV-MV Cookbook Reference: Section A2.5.1",
            font_size=12, color=MUTED)

add_footer(sl, 7, TOTAL)

# ====================================================================
#  SLIDE 8 — Coefficient & Residual QA/QC
# ====================================================================
sl = prs.slides.add_slide(BLANK_LAYOUT)
set_slide_bg(sl)
add_textbox(sl, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
            "Analytics \u2014 Coefficient & Residual QA/QC",
            font_size=28, bold=True, color=GREEN)
add_divider(sl, Inches(0.5), Inches(1.0), Inches(8.5), color=GREEN)

# Left column
add_textbox(sl, Inches(0.5), Inches(1.3), Inches(5), Inches(0.5),
            "Coefficient-Level", font_size=20, bold=True, color=ORANGE)
add_bullet_frame(sl, Inches(0.5), Inches(1.9), Inches(5.5), Inches(3),
                 [
                     "t-statistic \u2265 2.0 for significance",
                     "p-value < 0.05 for each predictor",
                     "Per-coefficient detail output:",
                     "  value, t-stat, p-value, significant flag",
                 ], font_size=14, spacing_pt=8)

# Right column
add_textbox(sl, Inches(7), Inches(1.3), Inches(5), Inches(0.5),
            "Residual Analysis", font_size=20, bold=True, color=RED)
add_bullet_frame(sl, Inches(7), Inches(1.9), Inches(5.5), Inches(3),
                 [
                     "Shapiro-Wilk normality test",
                     "Outlier detection (> 3\u03c3)",
                     "Reports outlier indices + count",
                     "Residual mean & std deviation",
                 ], font_size=14, spacing_pt=8)

# FSU box
shape = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                              Inches(2), Inches(5.0), Inches(9), Inches(1.8))
shape.fill.solid()
shape.fill.fore_color.rgb = BG_CARD
shape.line.color.rgb = ACCENT
shape.line.width = Pt(2)
add_textbox(sl, Inches(2), Inches(5.1), Inches(9), Inches(0.5),
            "Fractional Savings Uncertainty (FSU)",
            font_size=18, bold=True, color=ACCENT, alignment=PP_ALIGN.CENTER)
add_textbox(sl, Inches(2), Inches(5.6), Inches(9), Inches(0.4),
            "ASHRAE Guideline 14 Annex B  |  Accounts for autocorrelation",
            font_size=13, color=MUTED, alignment=PP_ALIGN.CENTER)
add_textbox(sl, Inches(2), Inches(6.1), Inches(9), Inches(0.4),
            "n\u2032 = n(1\u2013\u03c1)/(1+\u03c1)    |    FSU < 50% at 68% confidence = PASS",
            font_size=13, color=GREEN, alignment=PP_ALIGN.CENTER,
            font_name="Consolas")

add_footer(sl, 8, TOTAL)

# ====================================================================
#  SLIDE 9 — Data Flow Pipeline
# ====================================================================
sl = prs.slides.add_slide(BLANK_LAYOUT)
set_slide_bg(sl)
add_textbox(sl, Inches(0), Inches(0.2), SW, Inches(0.7),
            "Analytics Agent \u2014 Data Flow",
            font_size=30, bold=True, color=GREEN, alignment=PP_ALIGN.CENTER)
add_divider(sl, Inches(3), Inches(0.9), Inches(7.3), color=GREEN)

steps = [
    ("Input\nValidation",    "baseline_data\ndependent_var\npredictors",    ACCENT),
    ("OLS\nRegression",      "statsmodels\nadd_constant\nfit()",            GREEN),
    ("Model\nQA/QC",         "R\u00b2, CV(RMSE)\nNDBE, F-stat\nautocorr",  ORANGE),
    ("Coeff &\nResiduals",   "t-stat, p-val\nShapiro-Wilk\noutliers",      RED),
    ("Post-Period\n& FSU",   "savings calc\nuncertainty\neffective_n",      ACCENT2),
]

bw, bh = 2.0, 2.2
gap = 0.35
total_w = len(steps) * bw + (len(steps) - 1) * gap
start_x = (13.333 - total_w) / 2

for i, (title, detail, col) in enumerate(steps):
    bx = start_x + i * (bw + gap)
    by = 2.5

    # Step number circle
    circ = sl.shapes.add_shape(MSO_SHAPE.OVAL,
        Inches(bx + bw/2 - 0.18), Inches(1.5), Inches(0.36), Inches(0.36))
    circ.fill.solid()
    circ.fill.fore_color.rgb = col
    circ.line.fill.background()
    tf = circ.text_frame
    p = tf.paragraphs[0]
    p.text = str(i + 1)
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = BG
    p.alignment = PP_ALIGN.CENTER

    # Box
    shape = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(bx), Inches(by), Inches(bw), Inches(bh))
    shape.fill.solid()
    shape.fill.fore_color.rgb = BG_CARD
    shape.line.color.rgb = col
    shape.line.width = Pt(2)

    # Title in box
    add_textbox(sl, Inches(bx + 0.05), Inches(by + 0.05), Inches(bw - 0.1), Inches(0.6),
                title, font_size=12, bold=True, color=col, alignment=PP_ALIGN.CENTER)
    # Detail
    add_textbox(sl, Inches(bx + 0.1), Inches(by + 0.75), Inches(bw - 0.2), Inches(1.3),
                detail, font_size=10, color=MUTED, alignment=PP_ALIGN.CENTER,
                font_name="Consolas")

    # Arrow to next
    if i < len(steps) - 1:
        arr_x = bx + bw + 0.02
        conn = sl.shapes.add_connector(1,
            Inches(arr_x), Inches(by + bh/2),
            Inches(arr_x + gap - 0.04), Inches(by + bh/2))
        conn.line.color.rgb = WHITE
        conn.line.width = Pt(1.5)

# Annotation
add_textbox(sl, Inches(0), Inches(5.3), SW, Inches(0.5),
            "Output: Structured JSON dict with model stats, QA/QC results, savings, uncertainty",
            font_size=14, color=ACCENT, alignment=PP_ALIGN.CENTER)
add_textbox(sl, Inches(0), Inches(5.8), SW, Inches(0.4),
            "Ready to feed into Documentation Agent or render in Streamlit UI",
            font_size=12, color=MUTED, alignment=PP_ALIGN.CENTER)

add_footer(sl, 9, TOTAL)

# ====================================================================
#  SLIDE 10 — Output Schema
# ====================================================================
sl = prs.slides.add_slide(BLANK_LAYOUT)
set_slide_bg(sl)
add_textbox(sl, Inches(0.5), Inches(0.3), Inches(10), Inches(0.7),
            "Analytics Agent \u2014 Output Schema",
            font_size=30, bold=True, color=GREEN)
add_divider(sl, Inches(0.5), Inches(1.0), Inches(7), color=GREEN)

# Code block background
code_shape = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(0.4), Inches(1.2), Inches(8), Inches(5.8))
code_shape.fill.solid()
code_shape.fill.fore_color.rgb = RGBColor(0x0F, 0x19, 0x23)
code_shape.line.color.rgb = ACCENT2
code_shape.line.width = Pt(1.5)

schema = (
    '{\n'
    '  "model_type": "OLS",\n'
    '  "r2": 0.986,   "adj_r2": 0.981,\n'
    '  "cvrmse_percent": 4.3,\n'
    '  "ndbe_percent": 0.002,\n'
    '  "f_statistic": 142.7,  "f_pvalue": 1.2e-08,\n'
    '  "autocorrelation_lag1": 0.12,\n'
    '  "qa_qc": {\n'
    '    "model_pass": true,\n'
    '    "model_level": { "r2": {"pass": true}, ... },\n'
    '    "coefficient_level": { "const": {...}, "temp": {...} }\n'
    '  },\n'
    '  "residual_analysis": {\n'
    '    "normality_shapiro": {"p_value": 0.74},\n'
    '    "n_outliers": 0\n'
    '  },\n'
    '  "post_period": {\n'
    '    "estimated_savings": 124403,\n'
    '    "uncertainty": {"fsu": 0.281, "fsu_pass": true}\n'
    '  }\n'
    '}'
)
add_textbox(sl, Inches(0.6), Inches(1.3), Inches(7.5), Inches(5.5),
            schema, font_size=12, color=GREEN, font_name="Consolas")

# Annotations on right
annotations = [
    (1.6, "Core model metrics"),
    (3.0, "15+ QA/QC checks"),
    (4.6, "Residual diagnostics"),
    (5.6, "Post-period savings & FSU"),
]
for top, label in annotations:
    add_textbox(sl, Inches(9.0), Inches(top), Inches(4), Inches(0.4),
                "\u2190  " + label, font_size=13, color=ACCENT)

add_footer(sl, 10, TOTAL)

# ====================================================================
#  SLIDE 11 — API & Integration
# ====================================================================
sl = prs.slides.add_slide(BLANK_LAYOUT)
set_slide_bg(sl)
add_textbox(sl, Inches(0.5), Inches(0.3), Inches(8), Inches(0.7),
            "API & Integration", font_size=32, bold=True, color=WHITE)
add_divider(sl, Inches(0.5), Inches(1.0), Inches(5))

# Endpoints header
hdr = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE,
    Inches(0.5), Inches(1.3), Inches(12), Inches(0.5))
hdr.fill.solid()
hdr.fill.fore_color.rgb = ACCENT2
hdr.line.fill.background()
add_textbox(sl, Inches(0.7), Inches(1.3), Inches(2), Inches(0.5),
            "Method", font_size=13, bold=True, color=WHITE)
add_textbox(sl, Inches(2.5), Inches(1.3), Inches(2), Inches(0.5),
            "Endpoint", font_size=13, bold=True, color=WHITE)
add_textbox(sl, Inches(5.5), Inches(1.3), Inches(7), Inches(0.5),
            "Description", font_size=13, bold=True, color=WHITE)

endpoints = [
    ("GET",  "/health", "Service health check"),
    ("POST", "/chat",   "Route to agents via Orchestrator"),
]
for i, (meth, ep, desc) in enumerate(endpoints):
    t = Inches(1.9 + i * 0.55)
    add_textbox(sl, Inches(0.7), t, Inches(1.5), Inches(0.4),
                meth, font_size=13, color=GREEN, bold=True, font_name="Consolas")
    add_textbox(sl, Inches(2.5), t, Inches(2.5), Inches(0.4),
                ep, font_size=13, color=WHITE, font_name="Consolas")
    add_textbox(sl, Inches(5.5), t, Inches(7), Inches(0.4),
                desc, font_size=12, color=MUTED)

# Payload example
pay_shape = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(0.4), Inches(3.5), Inches(6), Inches(3.5))
pay_shape.fill.solid()
pay_shape.fill.fore_color.rgb = RGBColor(0x0F, 0x19, 0x23)
pay_shape.line.color.rgb = ACCENT2
pay_shape.line.width = Pt(1.5)

add_textbox(sl, Inches(0.6), Inches(3.5), Inches(5), Inches(0.4),
            "POST /chat \u2014 Example Payload", font_size=13, bold=True, color=ACCENT)
payload = (
    '{\n'
    '  "message": "Run regression",\n'
    '  "context": {\n'
    '    "dependent_var": "energy",\n'
    '    "predictors": ["temp","hours"],\n'
    '    "baseline_data": [{...}, ...]\n'
    '  }\n'
    '}'
)
add_textbox(sl, Inches(0.6), Inches(3.9), Inches(5.5), Inches(2.8),
            payload, font_size=12, color=WHITE, font_name="Consolas")

# Integration points
add_textbox(sl, Inches(7.5), Inches(3.5), Inches(5), Inches(0.5),
            "Integration Points", font_size=18, bold=True, color=ACCENT)
add_bullet_frame(sl, Inches(7.5), Inches(4.1), Inches(5), Inches(3),
                 [
                     "Streamlit UI (ui/app.py)",
                     "Swagger at /docs",
                     "Any REST client / SDK",
                     "CI/CD pipeline hooks",
                     "LLM orchestration layer",
                 ], font_size=13, spacing_pt=8)

add_footer(sl, 11, TOTAL)

# ====================================================================
#  SLIDE 12 — Roadmap
# ====================================================================
sl = prs.slides.add_slide(BLANK_LAYOUT)
set_slide_bg(sl)
add_textbox(sl, Inches(0), Inches(0.2), SW, Inches(0.7),
            "Roadmap & Next Steps", font_size=32, bold=True, color=WHITE,
            alignment=PP_ALIGN.CENTER)
add_divider(sl, Inches(3), Inches(0.9), Inches(7.3))

roadmap = [
    ("Now", [
        "OLS baseline regression",
        "15+ QA/QC checks",
        "FSU calculation",
        "Streamlit UI",
    ], GREEN),
    ("Next", [
        "Change-point models (3P/4P/5P)",
        "Multi-fuel support",
        "Interval data (hourly)",
        "Enhanced visualizations",
    ], ORANGE),
    ("Future", [
        "LLM-powered orchestration",
        "Auto-report generation",
        "Portfolio-level analytics",
        "ITC/ITV integration",
    ], RED),
]

for i, (phase, items, col) in enumerate(roadmap):
    bx = 0.7 + i * 4.2
    by = 1.5

    shape = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(bx), Inches(by), Inches(3.6), Inches(4.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = BG_CARD
    shape.line.color.rgb = col
    shape.line.width = Pt(2)

    add_textbox(sl, Inches(bx), Inches(by + 0.1), Inches(3.6), Inches(0.5),
                phase, font_size=20, bold=True, color=col, alignment=PP_ALIGN.CENTER)

    add_bullet_frame(sl, Inches(bx + 0.15), Inches(by + 0.7), Inches(3.3), Inches(3.5),
                     items, font_size=12, spacing_pt=8, color=WHITE)

    # Arrow to next phase
    if i < 2:
        arr_x = bx + 3.65
        conn = sl.shapes.add_connector(1,
            Inches(arr_x), Inches(by + 2.25),
            Inches(arr_x + 0.5), Inches(by + 2.25))
        conn.line.color.rgb = WHITE
        conn.line.width = Pt(1.5)

add_textbox(sl, Inches(0), Inches(6.5), SW, Inches(0.5),
            "github.com/sgfernandes/AIM-V",
            font_size=16, color=ACCENT, alignment=PP_ALIGN.CENTER,
            font_name="Consolas")

add_footer(sl, 12, TOTAL)

# ====================================================================
#  SAVE
# ====================================================================
out_path = "docs/AIM-V_Slides.pptx"
prs.save(out_path)
print(f"PowerPoint deck generated: {out_path}")
print(f"  {TOTAL} slides, 16:9 widescreen format")
